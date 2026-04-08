#!/usr/bin/env python3
"""Update MBR slides and create cheat sheet."""

from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn
from lxml import etree
import copy
import os


SECTIONS = {
    "RTAT Strategy": {
        "old": [
            "▪  REPAIR Framework:",
            "▪      Rank  |  Establish  |  Process  |  Amplify  |  Integrate  |  Redirect",
            "▪  10D target — once met, move to 8D",
            "▪  Selecting new focus ASCs for handed-off accounts",
            "▪  Engage: Doctor Maint (37.0D), Keisers (30.9D)",
            "▪  Urban: hold 10D, push to 8D by Q3",
            "▪  Rural: 11D by Q2, 10D by Q3",
        ],
        "new": [
            "▪  REPAIR Framework:",
            "▪      Rank  |  Establish  |  Process  |  Amplify  |  Integrate  |  Redirect",
            "▪  Shifting to rural-heavy focus — Urban hit 9.7D (below 10D target)",
            "▪  Continue holding Urban below 10D while driving Rural to 10D",
            "▪  Once 10D met across both markets, target moves to 8D before peak season",
            "▪  Dynamic Top 10: rank ASCs by volume × RTAT gap, apply REPAIR,",
            "▪  cycle list as accounts stabilize and new ones surface",
        ],
    },
    "Pending Review": {
        "old": [
            "▪  Urban: 93 top-10 pending, 10.6D avg delay",
            "▪  Rural: 53 top-10 pending, 17.4D avg delay",
            "▪  Appliance Medic: 59.8D avg delay",
            "▪  TV Specialty Shop: 24.8D avg delay",
        ],
        "new": [
            "▪  Started at ~16% of pending at 30+ days when we took over",
            "▪  Now at ~9.4% combined (Graham 7.6%, Brown 11.4%)",
            "▪  Total pending: 1,888 → 1,570 (W8→W12)",
            "▪  Avg pending age: 12.9D combined, trending down",
            "▪  60+ day pending: 20 total — pushing toward zero",
        ],
    },
    "Pending Strategy": {
        "old": [
            "▪  Appliance Medic 59.8D — immediate escalation",
            "▪  A-OK: 16 pending — accelerate scheduling",
            "▪  Target: top-10 avg delay below 10D",
        ],
        "new": [
            "▪  Pending managed holistically — not split by Urban/Rural",
            "▪  30+ day pending: ~16% → now ~9.4% (goal: 7%)",
            "▪  Push 30+ toward zero — remaining = backorder parts only",
            "▪  Once 30+ controlled, shift focus to 25-30 day bucket",
            "▪  Daily tracking via dashboard linked to Google Sheets",
            "▪  Avg pending days tracked WoW — target continuous reduction",
            "▪  Top 10 ASC avg pending age target: below 10 days",
        ],
    },
}


def set_shape_text(shape, new_text):
    """Replace ALL text in a shape, preserving formatting of the first run."""
    tf = shape.text_frame
    # Get first paragraph, first run
    para = tf.paragraphs[0]
    runs = para.runs
    if runs:
        # Keep first run's formatting, set its text, remove other runs
        runs[0].text = new_text
        # Remove extra runs via XML
        a_p = para._p
        r_elements = a_p.findall(qn('a:r'))
        for r in r_elements[1:]:
            a_p.remove(r)
    else:
        para.text = new_text


def make_shape_copy(slide, src_shape, new_text, new_y):
    """Deep-copy a shape XML element, change its text and y position, add to slide."""
    spTree = slide.shapes._spTree
    new_sp = copy.deepcopy(src_shape._element)

    # Change y position
    xfrm = new_sp.find('.//' + qn('a:xfrm'))
    if xfrm is not None:
        off = xfrm.find(qn('a:off'))
        if off is not None:
            off.set('y', str(new_y))

    # Change text: find all a:r elements, keep first, remove rest, set text on first
    all_runs = list(new_sp.iter(qn('a:r')))
    if all_runs:
        # Set text on first run
        t_elem = all_runs[0].find(qn('a:t'))
        if t_elem is not None:
            t_elem.text = new_text
        # Remove extra runs from their parent paragraphs
        for extra_r in all_runs[1:]:
            extra_r.getparent().remove(extra_r)

    spTree.append(new_sp)


def remove_shape(slide, shape):
    """Remove shape from slide."""
    shape._element.getparent().remove(shape._element)


def find_shapes_by_texts(slide, texts):
    """Find shapes matching a list of texts, in order."""
    text_to_shape = {}
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t in texts:
                text_to_shape[t] = shape
    return [text_to_shape.get(t) for t in texts]


def update_section(slide, section_name, section_def):
    """Replace bullets in a section."""
    old_texts = section_def["old"]
    new_texts = section_def["new"]

    old_shapes = find_shapes_by_texts(slide, old_texts)
    found = [s for s in old_shapes if s is not None]

    if not found:
        print(f"  WARNING: No shapes found for {section_name}")
        return

    # Use first found shape as template for position/formatting
    template = found[0]
    start_y = template.top
    line_spacing = Emu(182880)  # spacing between bullet lines

    # Remove ALL old shapes
    for shape in found:
        remove_shape(slide, shape)

    # Add new shapes
    for i, new_text in enumerate(new_texts):
        y = start_y + (i * line_spacing)
        make_shape_copy(slide, template, new_text, y)

    print(f"  {section_name}: {len(found)} → {len(new_texts)} bullets")


def update_pptx(filename):
    """Update a single PowerPoint file."""
    print(f"\nUpdating {filename}...")
    prs = Presentation(filename)
    slide = prs.slides[0]

    for name, sdef in SECTIONS.items():
        update_section(slide, name, sdef)

    prs.save(filename)
    print(f"  Saved {filename}")


def build_cheat_sheet():
    """Create the manager cheat sheet Word document."""
    from docx import Document
    from docx.enum.table import WD_TABLE_ALIGNMENT

    doc = Document()
    doc.add_heading('RTAT & Pending — Manager Coverage Cheat Sheet', level=1)
    doc.add_paragraph('Field Service | Urban & Rural Markets | April 2026')
    doc.add_paragraph('')

    # 1. REPAIR Framework
    doc.add_heading('1. REPAIR Framework', level=2)
    doc.add_paragraph(
        'REPAIR is the repeatable framework applied to every underperforming ASC. '
        'It cycles continuously as accounts improve and new ones surface on the Top 10 list.'
    )
    for label, desc in [
        ('R — Rank', 'Identify worst-performing ASCs using Impact Score = (ASC RTAT − 10 day goal) × volume. Highest score = most damage to the network. This decides who gets attention first.'),
        ('E — Establish', 'Set specific RTAT targets per ASC based on their current performance and gap from goal. Each ASC gets a clear number to hit.'),
        ('P — Process', 'Standardize claim filing — ensure ASCs are closing jobs promptly, not sitting on completed repairs. Process breakdowns are the #1 reason for inflated RTAT.'),
        ('A — Amplify', "Live Google Sheets dashboard gives visibility into every ASC's performance, trending, and pending in real time. No waiting for monthly reports."),
        ('I — Integrate', 'Cross-ASM collaboration — Brown and Graham share insights, coordinate on ASCs near territory borders, and align on strategy.'),
        ('R — Redirect', "Move volume away from ASCs that aren't improving to proven partners who are performing. Last resort, but necessary when engagement fails."),
    ]:
        p = doc.add_paragraph()
        p.add_run(label + ': ').bold = True
        p.add_run(desc)
    doc.add_paragraph('')

    # 2. RTAT Strategy
    doc.add_heading('2. RTAT Strategy — How It Works', level=2)
    doc.add_paragraph(
        'We maintain a dynamic Top 10 focus list for Urban and Rural. ASCs are ranked by '
        'Impact Score — highest volume × worst RTAT gap gets attention first. As an ASC '
        'improves and stabilizes, it falls off the list and the next worst ASC moves up.'
    )
    doc.add_paragraph(
        'This is not a fixed list. It continuously adjusts based on the most recent 4-week data. '
        'The strategy is to always be working the accounts causing the most damage to the network right now.'
    )
    doc.add_heading('Current State:', level=3)
    for b in [
        'Urban hit 9.7D — below the 10D target. Hold and continue pushing toward 8D.',
        'Rural at 12.8D, trending to 11.4D — this is now the primary focus area.',
        'Once 10D is met across both markets, target moves to 8D before peak season.',
        '7 ASCs have been stabilized and handed off to other managers — proof that REPAIR works.',
    ]:
        doc.add_paragraph(b, style='List Bullet')

    doc.add_heading('How REPAIR Applies (Example):', level=3)
    doc.add_paragraph(
        'An ASC averaging 25D RTAT on 200 jobs has an Impact Score of 3,000 '
        '(15 excess days × 200 jobs). That ASC gets prioritized on the Top 10 list. '
        'We establish a target, work the process issues, track weekly on the dashboard. '
        'When they improve, the next ASC in line gets the same treatment. The list is always moving.'
    )
    doc.add_paragraph('')

    # 3. Pending Strategy
    doc.add_heading('3. Pending Strategy — How It Works', level=2)
    doc.add_paragraph(
        'Pending is managed holistically across the entire book — not split by Urban/Rural. '
        'All open jobs are tracked daily.'
    )
    doc.add_heading('Key Metric: % of Pending at 30+ Days', level=3)
    for b in [
        'Started at ~16% when we took over.',
        'Currently at ~9.4% combined (Graham 7.6%, Brown 11.4%).',
        'Goal: 7% (configured in dashboard).',
        'End state: the only jobs at 30+ days should be backorder parts.',
    ]:
        doc.add_paragraph(b, style='List Bullet')

    doc.add_heading('Process:', level=3)
    for b in [
        'Pending data entered daily into Google Sheets.',
        'Custom dashboard tracks: total pending, 30+ day count/%, 60+ day count/%, average age, and week-over-week deltas.',
        'ASMs review pending on Monday and Thursday (per task schedule).',
        'Top 10 ASC avg pending age target: below 10 days.',
        'Once 30+ days is under control, shift focus to the 25-30 day bucket and keep driving the aging curve down.',
    ]:
        doc.add_paragraph(b, style='List Bullet')

    doc.add_heading('March Pending Recap (W8 → W12):', level=3)
    table = doc.add_table(rows=6, cols=4)
    table.style = 'Light Grid Accent 1'
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(['Metric', 'Graham', 'Brown', 'Combined']):
        cell = table.rows[0].cells[i]
        cell.text = h
        for r in cell.paragraphs[0].runs:
            r.bold = True
    for ri, rd in enumerate([
        ['Total Pending', '834', '736', '1,570'],
        ['30+ Days', '63 (7.6%)', '84 (11.4%)', '147 (~9.4%)'],
        ['60+ Days', '8 (1.0%)', '12 (1.6%)', '20 (~1.3%)'],
        ['Avg Age', '12.1D', '13.7D', '~12.9D'],
        ['WoW 30+ Trend', '\u2193 2.8pts', '\u2193 0.8pts', 'Improving'],
    ]):
        for ci, val in enumerate(rd):
            table.rows[ri + 1].cells[ci].text = val
    doc.add_paragraph('')

    # 4. Key Links & Schedule
    doc.add_heading('4. Key Links & Task Schedule', level=2)
    doc.add_heading('Task Schedule:', level=3)
    for s in [
        'Monday / Thursday: Pending Work, Top 10 Urban, Top 10 Rural',
        'Tuesday: Coverage Gaps',
        'Wednesday / Thursday: Top 30 Markets, Sears A&E Gaps',
    ]:
        doc.add_paragraph(s, style='List Bullet')
    doc.add_heading('Key Links:', level=3)
    for l in [
        'Pending (Graham): Google Sheet — "Dan & Trey" pending tracker',
        'Pending (Brown): Google Sheet — "Dan & Trey" pending tracker',
        'RTAT Dashboard: Google Sheet with per-ASM Urban/Rural tabs',
        'Coverage Gaps: Network Dashboard → TV/MS Coverage Gaps page',
    ]:
        doc.add_paragraph(l, style='List Bullet')
    doc.add_paragraph('')

    # 5. Quick Reference
    doc.add_heading('5. Quick Reference — Key Numbers', level=2)
    table2 = doc.add_table(rows=7, cols=3)
    table2.style = 'Light Grid Accent 1'
    table2.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, h in enumerate(['Metric', 'Current', 'Target']):
        cell = table2.rows[0].cells[i]
        cell.text = h
        for r in cell.paragraphs[0].runs:
            r.bold = True
    for ri, rd in enumerate([
        ['Urban RTAT', '9.7D', '10D (then 8D)'],
        ['Rural RTAT', '12.8D (trending 11.4D)', '10D (then 8D)'],
        ['Pending 30+ %', '~9.4%', '7%'],
        ['Pending Avg Age', '~12.9D', 'Continuous reduction'],
        ['Top 10 ASC Avg Pending Age', 'Varies', '< 10 days'],
        ['8D RTAT Target', 'After 10D met', 'Before peak season'],
    ]):
        for ci, val in enumerate(rd):
            table2.rows[ri + 1].cells[ci].text = val

    outpath = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'Manager_Cheat_Sheet.docx')
    doc.save(outpath)
    print(f"\nCheat sheet saved: {outpath}")


if __name__ == '__main__':
    os.chdir(os.path.dirname(os.path.abspath(__file__)))
    update_pptx('MBR_VersionA_Template.pptx')
    update_pptx('MBR_VersionB_Story.pptx')
    build_cheat_sheet()
    print("\nDone!")
